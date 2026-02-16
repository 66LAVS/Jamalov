from pathlib import Path

from docx import Document
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


DOCX_FILE = Path("Курсовая_Обзор_компьютерных_игр_Жамалов_Салман_Имаммединович.docx")
EXAMPLE_1 = Path("Пример/Пример ИП_1 .docx")
EXAMPLE_2 = Path("Пример/Пример ИП_2.docx")
PPTX_FILE = Path("Презентация_Обзор_компьютерных_игр_Жамалов_Салман_Имаммединович.pptx")
ASSETS_DIR = Path("ppt_assets")


def docx_summary(path: Path):
    d = Document(str(path))
    s = d.sections[0]
    style = d.styles["Normal"]
    size_pt = style.font.size.pt if style.font.size is not None else None
    return {
        "file": str(path),
        "paragraphs": len(d.paragraphs),
        "margins_cm": (
            round(s.left_margin.cm, 2),
            round(s.right_margin.cm, 2),
            round(s.top_margin.cm, 2),
            round(s.bottom_margin.cm, 2),
        ),
        "font": style.font.name,
        "font_size_pt": size_pt,
        "line_spacing": style.paragraph_format.line_spacing,
        "different_first_page": s.different_first_page_header_footer,
    }


def make_image(path: Path, title: str, subtitle: str, color=(28, 45, 86)):
    img = Image.new("RGB", (1280, 720), color)
    draw = ImageDraw.Draw(img)
    try:
        font_title = ImageFont.truetype("arial.ttf", 64)
        font_sub = ImageFont.truetype("arial.ttf", 34)
    except OSError:
        font_title = ImageFont.load_default()
        font_sub = ImageFont.load_default()

    draw.rectangle((60, 60, 1220, 660), outline=(255, 255, 255), width=4)
    draw.text((110, 220), title, fill=(255, 255, 255), font=font_title)
    draw.text((110, 330), subtitle, fill=(225, 225, 225), font=font_sub)
    img.save(path)


def set_title(slide, text):
    title = slide.shapes.title
    title.text = text
    p = title.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.name = "Times New Roman"
        run.font.bold = True
        run.font.size = Pt(36)
        run.font.color.rgb = RGBColor(24, 24, 24)


def set_body(tf, lines):
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = "Times New Roman"
            run.font.size = Pt(22)
            run.font.color.rgb = RGBColor(30, 30, 30)


def create_presentation():
    ASSETS_DIR.mkdir(exist_ok=True)

    image_specs = [
        ("01_title.png", "Обзор компьютерных игр", "Индивидуальный проект"),
        ("02_relevance.png", "Актуальность", "Игровая индустрия как часть цифровой экономики"),
        ("03_goal.png", "Цель и задачи", "Структура исследовательской работы"),
        ("04_object.png", "Объект и предмет", "Методы исследования"),
        ("05_history.png", "История игр", "От аркад до современных AAA"),
        ("06_genres.png", "Жанры", "Экшен, RPG, стратегии, симуляторы"),
        ("07_platforms.png", "Платформы", "ПК, консоли, мобильные устройства"),
        ("08_society.png", "Влияние на общество", "Плюсы и возможные риски"),
        ("09_games.png", "Практический обзор", "Counter-Strike 2, Dota 2, Minecraft"),
        ("10_compare.png", "Сравнительный анализ", "Графика, сюжет, геймплей"),
        ("11_choice.png", "Выбор лучшего продукта", "Обоснование выбора Minecraft"),
        ("12_feedback.png", "Отзывы пользователей", "Оценки, рейтинги, ожидания"),
        ("13_conclusion.png", "Выводы", "Итоги исследования"),
        ("14_thanks.png", "Спасибо за внимание", "Вопросы?"),
    ]

    for name, title, subtitle in image_specs:
        make_image(ASSETS_DIR / name, title, subtitle)

    prs = Presentation()

    slides_data = [
        ("Тема проекта", ["Обзор компьютерных игр", "Жамалов Салман Имаммединович", "Нижневартовск, 2026"], "01_title.png"),
        ("Актуальность темы", ["Игры — значимая часть цифровой индустрии", "Влияние на образование, культуру, экономику", "Рост киберспорта и онлайн-сервисов"], "02_relevance.png"),
        ("Цель и задачи", ["Цель: провести обзор компьютерных игр", "Задачи: история, жанры, платформы", "Сравнение и выводы по качеству"], "03_goal.png"),
        ("Объект, предмет, методы", ["Объект: компьютерные игры", "Предмет: жанровые и технологические особенности", "Методы: анализ, сравнение, обзор оценок"], "04_object.png"),
        ("1.1 Эволюция игр", ["Аркадные автоматы и первые консоли", "Переход к 3D и сетевым режимам", "Современная модель: «игра как сервис»"], "05_history.png"),
        ("1.2 Жанры игр", ["Экшен — динамика и реакция", "RPG — развитие персонажа и сюжет", "Стратегии/симуляторы — планирование"], "06_genres.png"),
        ("1.3 Платформы и технологии", ["ПК: гибкость и модификации", "Консоли: стабильность и оптимизация", "Unity и Unreal Engine"], "07_platforms.png"),
        ("1.4 Влияние игр на общество", ["Плюсы: развитие навыков и коммуникации", "Риски: зависимость, монетизация, токсичность", "Нужна цифровая грамотность"], "08_society.png"),
        ("2.1 Выбранные игры", ["Counter-Strike 2", "Dota 2", "Minecraft, Cyberpunk 2077, Forza Horizon 5"], "09_games.png"),
        ("2.2 Сравнительный анализ", ["Критерии: графика, геймплей, сюжет", "Оценка технической стабильности", "Уровень входа для новичков"], "10_compare.png"),
        ("2.3 Обоснование выбора", ["Наиболее сбалансированный продукт: Minecraft", "Доступность и образовательный потенциал", "Сильное сообщество и обновления"], "11_choice.png"),
        ("2.4 Пользовательский опыт", ["Отзывы зависят от ожиданий аудитории", "Важны: баланс, честный матчмейкинг", "Прозрачные обновления и поддержка"], "12_feedback.png"),
        ("Заключение", ["Цель проекта достигнута", "Проведен теоретический и практический обзор", "Определены перспективы развития индустрии"], "13_conclusion.png"),
        ("Спасибо за внимание", ["Готов ответить на вопросы", "Информационная база: 15 источников", "Презентация по материалам DOCX"], "14_thanks.png"),
    ]

    for title, bullets, img_name in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_title(slide, title)
        set_body(slide.placeholders[1].text_frame, bullets)
        slide.shapes.add_picture(str(ASSETS_DIR / img_name), Inches(7.0), Inches(1.3), width=Inches(6.0), height=Inches(3.4))

    prs.save(str(PPTX_FILE))


def main():
    target = docx_summary(DOCX_FILE)
    ex1 = docx_summary(EXAMPLE_1)
    ex2 = docx_summary(EXAMPLE_2)

    print("DOCX_CHECK_START")
    for item in (target, ex1, ex2):
        print(item)
    print("DOCX_CHECK_END")

    create_presentation()
    print(f"PPTX_CREATED: {PPTX_FILE}")


if __name__ == "__main__":
    main()
